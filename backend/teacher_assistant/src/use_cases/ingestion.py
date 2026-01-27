import os
import pdfplumber
from pptx import Presentation
from docx import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from ..infrastructure.database import VectorDatabase
from ..infrastructure.ollama_client import OllamaClient
import pandas as pd
from concurrent.futures import ThreadPoolExecutor, as_completed
import multiprocessing

class IngestionService:
    # GLOBAL PROGRESS TRACKER (CourseID -> Status)
    _progress_map = {}

    def __init__(self, db: VectorDatabase, llm: OllamaClient, rag_service=None, course_id: str = "default"):
        self.db = db
        self.llm = llm
        self.rag_service = rag_service
        self.course_id = course_id
        # Initialize progress for this session
        IngestionService._progress_map[self.course_id] = {"status": "starting", "progress": 0, "current_file": ""}
        
        # OPTIMIZED: Larger chunks for better context preservation
        self.splitter = RecursiveCharacterTextSplitter(
            chunk_size=800,      # Full paragraphs
            chunk_overlap=100    # Better continuity
        )

    def process_directory(self, directory: str):
        print(f"💎 KNOWLEDGE INGESTION STARTING for {self.course_id}...")
        IngestionService._progress_map[self.course_id] = {"status": "scanning", "progress": 5, "current_file": ""}
        
        all_files = []
        for root, _, files in os.walk(directory):
            for file in files:
                if file.lower().endswith(('.pdf', '.pptx', '.docx', '.txt', '.xlsx')):
                   all_files.append(os.path.join(root, file))

        if not all_files:
             print("⚠️ No supported files found!")
             IngestionService._progress_map[self.course_id] = {"status": "ready", "progress": 100, "current_file": ""}
             return

        print(f"🚀 PARALLEL PARSING: Processing {len(all_files)} files on {multiprocessing.cpu_count()} cores...")
        IngestionService._progress_map[self.course_id] = {"status": "parsing", "progress": 10, "current_file": f"Parallel Batch ({len(all_files)} docs)"}

        all_chunks = [] # Fix: Initialize collector
        # BEST PARALLEL COMPUTING: ThreadPool for I/O bound parsing
        with ThreadPoolExecutor(max_workers=min(len(all_files), 10)) as executor:
            future_to_file = {executor.submit(self._parse_file, path): path for path in all_files}
            
            completed_count = 0
            for future in as_completed(future_to_file):
                path = future_to_file[future]
                try:
                    chunks = future.result()
                    all_chunks.extend(chunks)
                    completed_count += 1
                    # Live Update
                    IngestionService._progress_map[self.course_id] = {
                        "status": "parsing", 
                        "progress": 10 + int((completed_count/len(all_files))*10), 
                        "current_file": os.path.basename(path)
                    }
                except Exception as exc:
                    print(f"❌ Error parsing {path}: {exc}")

        print(f"🧠 Embedding {len(all_chunks)} chunks on GPU...")
        IngestionService._progress_map[self.course_id] = {"status": "embedding", "progress": 20, "current_file": "GPU Indexing Core"}
        
        batch_size = 50 
        total_batches = (len(all_chunks) + batch_size - 1) // batch_size
        
        for i in range(0, len(all_chunks), batch_size):
            batch_idx = i // batch_size
            progress_val = 20 + int((batch_idx / total_batches) * 60)
            IngestionService._progress_map[self.course_id] = {"status": "embedding", "progress": progress_val, "current_file": f"Batch {batch_idx+1}/{total_batches}"}
            
            batch = all_chunks[i : i + batch_size]
            texts = [c['content'] for c in batch]
            try:
                vectors = self.llm.get_embeddings_batch(texts)
                for j, v in enumerate(vectors):
                    batch[j]['vector'] = v
            except Exception as e:
                print(f"⚠️ Batch failed: {e}")
                for c in batch:
                    c['vector'] = self.llm.get_embedding(c['content'])
            
        IngestionService._progress_map[self.course_id] = {"status": "saving", "progress": 85, "current_file": "Vector Space"}
        self.db.insert_chunks(all_chunks)
        
        # TRIGGER SYNTHETIC WARMING
        if self.rag_service:
            IngestionService._progress_map[self.course_id] = {"status": "caching", "progress": 90, "current_file": "Smart Warm-up"}
            self._warm_up_cache(all_chunks)
            
        IngestionService._progress_map[self.course_id] = {"status": "ready", "progress": 100, "current_file": ""}
        print(f"✅ Indexed {len(all_chunks)} chunks for {self.course_id}")

    def _warm_up_cache(self, chunks):
        print("\n🚀 STARTING PERFORMANCE PRE-FETCH...")
        import random
        
        # Pick 10 random chunks for critical pre-fetch
        samples = random.sample(chunks, min(len(chunks), 10))
        generated_qs = []
        
        for i, chunk in enumerate(samples):
            text = chunk['content'][:1000]
            prompt = (
                f"TEXT: {text}\n\n"
                "TASK: Generate 3 specific student questions about this text.\n"
                "FORMAT: Questions only, one per line."
            )
            try:
                resp = self.llm.chat("You are a knowledge-extraction tool.", prompt)
                qs = [q.strip() for q in resp.split('\n') if '?' in q]
                generated_qs.extend(qs)
            except Exception:
                pass
                
        print(f"🔄 Hydrating {len(generated_qs)} predicted entries...")
        for q in generated_qs:
            if len(q) < 5: continue
            try:
                self.rag_service.answer_question(q)
            except Exception:
                pass
                
        print(f"✅ PRE-FETCH COMPLETE: {len(generated_qs)} entries stabilized.")

    def _parse_file(self, path: str):
        print(f"DEBUG: Parsing file: {path}")
        ext = os.path.splitext(path)[1].lower()
        fname = os.path.basename(path)
        # Clean filename for embedding (remove extension and numbers)
        import re
        clean_name = fname
        for e in [".pptx", ".docx", ".pdf", ".xlsx"]:
            clean_name = clean_name.replace(e, "")
        clean_name = re.sub(r'^\d+\s*[-–]\s*', '', clean_name).strip()
        
        blocks = []
        try:
            if ext == ".pptx":
                prs = Presentation(path)
                for i, slide in enumerate(prs.slides):
                    # Get slide title if available
                    title = ""
                    for shape in slide.shapes:
                        if shape.has_text_frame and shape.text_frame.paragraphs:
                            title = shape.text_frame.paragraphs[0].text[:50]
                            break
                    text = " ".join([s.text for s in slide.shapes if hasattr(s, "text")])
                    if text.strip():
                        loc = f"Slide {i+1}" + (f": {title}" if title else "")
                        blocks.append({"text": text, "loc": loc})
            elif ext == ".pdf":
                with pdfplumber.open(path) as pdf:
                    for i, page in enumerate(pdf.pages):
                        text = page.extract_text()
                        if text:
                            blocks.append({"text": text, "loc": f"Page {i+1}"})
            elif ext == ".docx":
                doc = Document(path)
                # Group paragraphs for better context
                current_text = ""
                for i, p in enumerate(doc.paragraphs):
                    if p.text.strip():
                        current_text += p.text + " "
                        if len(current_text) > 500:  # Group small paragraphs
                            blocks.append({"text": current_text, "loc": f"Section {len(blocks)+1}"})
                            current_text = ""
                if current_text.strip():
                    blocks.append({"text": current_text, "loc": f"Section {len(blocks)+1}"})
            elif ext == ".txt":
                with open(path, "r", encoding="utf-8") as f:
                    text = f.read()
                    if text.strip():
                        blocks.append({"text": text, "loc": "Full Document"})
            elif ext == ".xlsx":
                # PANDAS MAGIC: Read all sheets
                xls = pd.ExcelFile(path)
                for sheet_name in xls.sheet_names:
                    df = pd.read_excel(xls, sheet_name=sheet_name)
                    # Convert to string representation
                    text = df.to_string(index=False)
                    if len(text) > 50: # Ignore empty sheets
                        blocks.append({"text": text, "loc": f"Sheet: {sheet_name}"})
                        print(f"📊 EXCEL DEBUG: Extracted {len(text)} chars from sheet '{sheet_name}'")
                    else:
                        print(f"⚠️ EXCEL DEBUG: Sheet '{sheet_name}' was empty or too short.")
        except Exception as e:
            print(f"⚠️ Error parsing {fname}: {e}")
            import traceback
            traceback.print_exc()

        final_chunks = []
        for b in blocks:
            splits = self.splitter.split_text(b['text'])
            for s in splits:
                # EMBED filename in content for FTS precision
                enriched_content = f"[{clean_name}] {s}"
                final_chunks.append({
                    "content": enriched_content,
                    "source": fname,
                    "location": b['loc']
                })
        return final_chunks
