
import requests
import time
import os
import pandas as pd
from docx import Document
from pptx import Presentation

# --- CONFIG ---
BASE_URL = "http://localhost:8000"
TEACHER_EMAIL = "teacher.stress@iitu.kz"
COURSE_ID = "stress_test_upload"

def create_dummy_files():
    print("📂 Generating dummy files for upload stress test...")
    os.makedirs("test_files", exist_ok=True)
    
    # 1. Excel (XLSX)
    df = pd.DataFrame({'Student': ['Alice', 'Bob', 'Charlie'], 'Grade': [85, 92, 78], 'Comments': ['Great', 'Excellent', 'Good']})
    df.to_excel("test_files/grades.xlsx", index=False)
    
    # 2. PowerPoint (PPTX)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Lecture 1: Introduction"
    slide.placeholders[1].text = "Welcome to the Stress Test Course."
    prs.save("test_files/lecture.pptx")
    
    # 3. Word (DOCX)
    doc = Document()
    doc.add_heading('Syllabus', 0)
    doc.add_paragraph('This is a test document covering all topics.')
    doc.save("test_files/syllabus.docx")
    
    # 4. Text (TXT)
    with open("test_files/notes.txt", "w") as f:
        f.write("Important exam dates: Midterm Oct 10, Final Dec 20.")
        
    print("✅ Files generated successfully.")

def login_teacher():
    # Register/Login
    print("\n🔑 Logging in Teacher...")
    try:
        requests.post(f"{BASE_URL}/api/auth/register", json={"email": TEACHER_EMAIL, "password": "pass", "name": "Stress Teacher"})
    except: pass
    
    resp = requests.post(f"{BASE_URL}/api/auth/login", json={"email": TEACHER_EMAIL, "password": "pass"})
    if resp.status_code == 200:
        return resp.json()["token"]
    print(f"❌ Login failed: {resp.text}")
    return None

def run_stress_test():
    create_dummy_files()
    token = login_teacher()
    if not token: return
    headers = {"Authorization": f"Bearer {token}"}
    
    # 1. Create Course
    print("\n📚 Creating Course Workspace...")
    requests.post(f"{BASE_URL}/api/courses", json={"id": COURSE_ID, "subject": "Stress Upload 101", "teacherName": "Mr. Stress", "teacherId": "t_stress"}, headers=headers)
    
    # 2. Upload Files
    print("\n🚀 Uploading Files (XLSX, PPTX, DOCX, TXT)...")
    files = [
        ('files', ('grades.xlsx', open('test_files/grades.xlsx', 'rb'), 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')),
        ('files', ('lecture.pptx', open('test_files/lecture.pptx', 'rb'), 'application/vnd.openxmlformats-officedocument.presentationml.presentation')),
        ('files', ('syllabus.docx', open('test_files/syllabus.docx', 'rb'), 'application/vnd.openxmlformats-officedocument.wordprocessingml.document')),
        ('files', ('notes.txt', open('test_files/notes.txt', 'rb'), 'text/plain'))
    ]
    
    ts = time.time()
    resp = requests.post(f"{BASE_URL}/api/upload", data={"course_id": COURSE_ID}, files=files, headers=headers)
    elapsed = time.time() - ts
    
    if resp.status_code == 200:
        print(f"✅ Upload Success in {elapsed:.2f}s: {resp.json()}")
    else:
        print(f"❌ Upload Failed: {resp.status_code} - {resp.text}")
        return

    # 3. Monitor Processing
    print("\n⏳ Monitoring Parallel Ingestion Status...")
    for _ in range(30):
        time.sleep(1)
        status = requests.get(f"{BASE_URL}/api/ingest/status/{COURSE_ID}").json()
        print(f"   Status: {status['status']} | Progress: {status['progress']}% | File: {status['current_file']}")
        if status['status'] == 'ready' or status['progress'] >= 100:
            print("✅ Ingestion COMPLETE!")
            break
            
    # 4. Verify Content
    print("\n🔍 Verifying Content (Chat Query)...")
    q_resp = requests.post(f"{BASE_URL}/api/chat", json={"message": "Who got Excellent grade?", "course_id": COURSE_ID}, headers=headers).json()
    print(f"   Q: Who got Excellent grade? (Expected: Bob)")
    print(f"   A: {q_resp['response']}")
    
    if "Bob" in q_resp['response'] or "92" in q_resp['response']:
        print("✅ SUCCESS: Excel data retrieved correctly!")
    else:
        print("⚠️ WARNING: Excel recall not as expected.")

if __name__ == "__main__":
    run_stress_test()
